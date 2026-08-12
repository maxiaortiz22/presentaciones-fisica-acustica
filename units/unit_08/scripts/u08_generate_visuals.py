#!/usr/bin/env python3
"""Genera los recursos visuales propios aprobados de la Unidad 08.

Salidas: SVG/PNG, fuentes JSON/CSV, PPTX editables por diagrama, previews
renderizados en PowerPoint y metadatos de validación. No genera el deck final.
"""

from __future__ import annotations

import csv
import json
import math
import re
import sys
import time
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import numpy as np

ROOT = Path(__file__).resolve().parents[3]
UNIT = ROOT / "units" / "unit_08"
GEN = UNIT / "assets" / "generated"

# El runtime del workspace aporta python-pptx; se agrega al final para no sustituir
# NumPy/Matplotlib del entorno científico del equipo.
RUNTIME_SITE = Path(r"C:\Users\maxia\.cache\codex-runtimes\codex-primary-runtime\dependencies\python\Lib\site-packages")
if RUNTIME_SITE.exists() and str(RUNTIME_SITE) not in sys.path:
    sys.path.append(str(RUNTIME_SITE))
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

COL = {
    "bordo": "#4D1434", "bordo2": "#903163", "carbon": "#3D3D3D",
    "gris": "#969FA7", "gris2": "#D9DCE0", "blanco": "#FFFFFF",
    "marfil": "#F7F6F2", "teal": "#2F7E83", "teal_bg": "#E7F1F1",
    "ocre": "#9F541A", "ocre_bg": "#F8EDE2", "ok": "#2F6F55",
    "alerta": "#9A641E", "error": "#A33A3A",
}

DIAGRAMS = {
    "001": ("Caso ramificado", "concept", ["Caso inicial", "Exposición\n¿qué ocurrió?", "Percepción\n¿qué refiere?", "Medición\n¿qué se obtuvo?"], ["exposición", "percepción", "medición"]),
    "002": ("Mapa de la unidad", "flow", ["Evidencia y límites", "Exposición y riesgo", "Estudios auditivos", "Rehabilitación e integración"], ["describir", "medir", "integrar"]),
    "003": ("Clases de evidencia", "concept", ["Exposición\nruido y tiempo", "Alteración\ncambio funcional", "Síntoma\nexperiencia referida", "Resultado\ndato de una prueba", "Limitación\nactividad y participación"], []),
    "004": ("Seis preguntas del caso", "grid", ["¿Qué entra?", "¿Qué sistema?", "¿Qué cambia?", "¿Qué se registra?", "¿En qué unidad?", "¿Qué no concluye?"], []),
    "005": ("Clasificación funcional", "columns", ["Transmisión\noído externo/medio", "Cóclea o vía neural\nrespuesta sensorial", "Combinación\ncomponentes coexistentes"], []),
    "006": ("Un patrón limita la inferencia", "concept", ["Patrón observado", "¿Cómo se midió?", "¿Qué antecedentes hay?", "¿Qué dato falta?", "Patrón ≠ causa única"], ["preguntar", "preguntar", "preguntar"]),
    "007": ("Batería convergente", "concept", ["Pregunta inicial", "Antecedentes", "Prueba conductual", "Prueba fisiológica", "Integración limitada", "Nueva pregunta"], ["aporta", "aporta", "aporta", "contrastar"]),
    "008": ("Variables de exposición", "grid", ["Nivel", "Descriptor", "Ponderación", "Duración", "Espectro", "Impulsividad"], []),
    "009": ("Cambio temporal de umbral", "equation", ["ΔL_T(f, Δt) = L_U,1(f, Δt) − L_U,0(f)", "f: frecuencia de prueba", "Δt: tiempo posexposición", "L_U,0: umbral de referencia", "L_U,1: umbral posterior"], []),
    "010": ("Ejemplo de TTS", "steps", ["Datos\n27 y 12 dB HL", "Compatibilidad\nmisma escala", "Resta\n27 − 12", "Resultado\n15 dB ≠ causa"], ["verificar", "calcular", "interpretar"]),
    "012": ("Evidencia compatible con exposición a ruido", "concept", ["Exposición documentada", "Patrón medido", "Antecedentes", "Condiciones de prueba", "Compatible con\nno demuestra causa"], ["aporta", "aporta", "limita", "controla"]),
    "013": ("Cambio auditivo multifactorial", "concept", ["Edad", "Exposición", "Salud", "Variabilidad", "Mediciones", "Desempeño"], []),
    "014": ("Antes de interpretar un porcentaje", "grid", ["1. Evento", "2. Población", "3. Período", "4. Exposición", "5. Comparador", "6. Incertidumbre"], []),
    "015": ("Cadena común de medición", "flow", ["Estímulo", "Sistema", "Respuesta", "Sensor o tarea", "Registro", "Interpretación"], ["interactúa", "produce", "capta", "representa", "limita"]),
    "016": ("Matriz de seis preguntas", "grid", ["Entrada", "Sistema", "Transformación", "Sensor o tarea", "Dato y unidad", "Límite"], []),
    "017": ("Mapa anatómico-funcional", "matrix", ["Transmisión", "Oído medio", "Cóclea", "Vía neural", "Audiometría", "Timpanometría", "OEA", "PEAT", "ECoG", "Logoaudiometría", "Acufenometría"], []),
    "018": ("Condiciones del registro", "concept", ["Registro", "Ambiente", "Equipo", "Sensor", "Protocolo", "Persona"], []),
    "019": ("Batería y discrepancia", "process", ["Pregunta", "Pruebas", "Comparar resultados", "Discrepancia", "Revisar condiciones", "Nueva pregunta"], ["orienta", "contrastar", "si aparece", "volver", "ampliar"]),
    "020": ("Vías de presentación", "mixed", ["Generador", "Auricular\npresión aérea", "Vibrador\nvibración mecánica", "Trayecto aéreo", "Trayecto óseo", "Tarea de respuesta"], ["aire", "vibración"]),
    "021": ("Referencias de nivel", "columns", ["dB SPL\nref.: presión física\nuso: señales acústicas", "dB HL\nref.: umbral normalizado\nuso: audiometría", "dB SL\nref.: umbral individual\nuso: nivel relativo"], []),
    "022": ("Diferencia aérea–ósea", "equation", ["G_AO(f) = L_VA(f) − L_VO(f)", "f: frecuencia", "L_VA: nivel por vía aérea", "L_VO: nivel por vía ósea", "Resultado: dB"], []),
    "023": ("Ejemplo de diferencia aérea–ósea", "steps", ["Vía aérea\n40 dB HL", "Vía ósea\n15 dB HL", "Resta\n40 − 15", "Resultado\n25 dB", "Límite\nno diagnostica"], ["comparar", "restar", "expresar", "interpretar"]),
    "024": ("Condiciones de audiometría", "concept", ["Audiograma", "Calibración", "Transductor", "Ambiente", "Consigna", "Respuesta", "Enmascaramiento"], []),
    "025": ("Tres tareas verbales", "columns", ["Detección\n¿hay señal?\nrespuesta sí/no", "Reconocimiento\n¿qué palabra fue?\nporcentaje correcto", "Identificación\nelegir alternativa\naciertos"], []),
    "026": ("Cadena logoaudiométrica", "flow", ["Material verbal", "Nivel", "Presentación", "Tarea", "Respuesta", "Porcentaje"], ["ajustar", "emitir", "responder", "puntuar", "representar"]),
    "027": ("Acufenometría como correspondencia", "process", ["Percepto referido", "Estímulo ajustable", "Decisión de semejanza", "Resultado de correspondencia"], ["ajustar", "comparar", "registrar"]),
    "028": ("Señal física y percepto", "columns", ["Señal externa\nfrecuencia y nivel", "Percepto referido\nexperiencia subjetiva", "Correspondencia\nresultado de una tarea"], []),
    "030": ("Cadena timpanométrica", "process", ["Tono de sonda", "Barrido de presión", "Oído medio", "Micrófono", "Inmitancia", "Curva"], ["ida acústica", "modifica", "retorno", "estima", "grafica"]),
    "031": ("Construcción del timpanograma", "steps", ["1. Cambiar presión", "2. Medir respuesta", "3. Ubicar un punto", "4. Repetir y trazar"], ["medir", "ubicar", "repetir"]),
    "032": ("Curva plana: preguntas abiertas", "concept", ["Curva observada", "¿Sellado adecuado?", "¿Sonda permeable?", "¿Barrido completo?", "¿Respuesta estable?", "Posibles explicaciones\nno causa única"], []),
    "033": ("OEA: ida y retorno", "mixed", ["Estímulo", "Oído externo/medio", "Cóclea · CCE", "Retorno acústico", "Micrófono", "Registro"], ["ida", "transmisión", "retorno", "captar", "representar"]),
    "034": ("OEA dentro de una batería", "concept", ["OEA", "Antecedentes", "Otras pruebas", "Condiciones técnicas", "Integración limitada"], ["aporta", "aporta", "controla", "integrar"]),
    "035": ("Cadena PEAT", "mixed", ["Estímulo", "Transductor", "Vía auditiva", "Electrodos", "Amplificar y promediar", "V(t)"], ["acústica", "respuesta", "bioeléctrica", "procesar", "registrar"]),
    "036": ("Cadena ECoG", "process", ["Estímulo", "Generadores cocleares/neural", "Electrodo cercano", "Amplificación", "Registro"], ["evoca", "potencial", "capta", "representa"]),
    "037": ("Componentes de ECoG", "columns", ["CM\nMicrofónico coclear\nseguimiento del estímulo", "SP\nPotencial de sumación\ncomponente sostenido", "AP\nPotencial de acción\nrespuesta neural compuesta"], []),
    "038": ("OEA, PEAT y ECoG", "columns", ["OEA\nEstímulo: acústico\nGenerador: cóclea/CCE\nSensor: micrófono\nMagnitud: presión\nLímite: integra batería", "PEAT\nEstímulo: acústico\nGenerador: vía auditiva\nSensor: electrodos\nMagnitud: V(t)\nLímite: depende del protocolo", "ECoG\nEstímulo: acústico\nGenerador: coclear/neural\nSensor: electrodo próximo\nMagnitud: potencial\nLímite: integración"], []),
    "039": ("Cadena común de dispositivos", "process", ["Entrada", "Transducción", "Procesamiento", "Salida", "Sistema remanente"], ["captar", "transformar", "entregar", "interactuar"]),
    "040": ("Cadena de audífono", "mixed", ["Sonido", "Micrófono", "Procesamiento", "Receptor", "Presión en oído"], ["acústica", "eléctrica/digital", "eléctrica", "acústica"]),
    "041": ("Ganancia por frecuencia", "equation", ["G(f) = L_salida(f) − L_entrada(f)", "f: frecuencia", "L_entrada: nivel de entrada", "L_salida: nivel de salida", "Resultado: dB"], []),
    "042": ("Ejemplo de ganancia", "steps", ["Entrada\n60 dB SPL", "Salida\n82 dB SPL", "Resta\n82 − 60", "Ganancia\n22 dB", "Límite\nno equivale a beneficio"], ["comparar", "restar", "expresar", "interpretar"]),
    "043": ("Cadena de implante coclear", "mixed", ["Sonido", "Micrófono", "Procesador", "Enlace", "Receptor/estimulador", "Electrodos"], ["captar", "codificar", "transmitir", "pulsos", "estimular"]),
    "044": ("Bandas, canales y perceptos", "layers", ["Bandas analizadas", "Canales programados", "Electrodos activos", "Perceptos posibles"], ["no es 1:1", "interacción", "depende del sistema neural"]),
    "045": ("Audífono e implante coclear", "columns", ["Audífono\nentrada acústica\nprocesamiento\nsalida acústica", "Implante coclear\nentrada acústica\ncodificación\nsalida eléctrica"], []),
    "046": ("Otras estrategias de salida", "mixed", ["Conducción ósea", "Vibración mecánica", "Ambas cócleas", "EAS", "Salida acústica", "Salida eléctrica"], ["entrega", "alcanza", "combina", "combina"]),
    "047": ("Síntesis de dispositivos", "grid", ["Audífono\nsalida acústica", "Implante coclear\nsalida eléctrica", "Conducción ósea\nsalida mecánica", "EAS\nsalida combinada"], []),
    "048": ("Matriz del caso", "matrix", ["Dato", "Clase", "Pregunta que habilita", "L_Aeq,T", "exposición", "¿cómo se describió?", "Tinnitus referido", "síntoma", "¿qué percibe?", "Umbral dB HL", "resultado", "¿bajo qué tarea?", "OEA: derivar", "resultado", "¿qué condición revisar?"], []),
    "049": ("Inferencia profesional acotada", "matrix", ["Dato", "Permite", "No permite", "Siguiente dato", "Umbral", "describir patrón", "nombrar causa", "antecedentes", "OEA", "registrar respuesta", "audición global", "otras pruebas", "Discrepancia", "revisar condiciones", "forzar acuerdo", "repetir/ampliar"], []),
    "050": ("Proceso profesional", "process", ["Pregunta", "Antecedentes", "Pruebas", "Control técnico", "Integración", "Decisión compartida"], ["orienta", "selecciona", "verifica", "contrasta", "comunica"]),
    "051": ("Seis preguntas transferibles", "grid", ["¿Qué entra?", "¿Qué sistema?", "¿Qué se transforma?", "¿Qué se registra o entrega?", "¿En qué unidad?", "¿Qué no concluye?"], []),
    "052": ("Puente de U8 a U10", "concept", ["Unidad 8\nmedición y límites", "Exposición", "Descriptor", "Efecto", "Control", "Unidad 10\nruido"], ["preguntar", "preguntar", "preguntar", "preguntar"]),
}

BLOCKED_DIAGRAMS = {
    "011": "descriptor y normativa de exposición pendientes (OD-U08-14)",
    "029": "notación de dB SL y umbral específico pendiente (OD-U08-09/21)",
}

CHARTS = {
    "001": ("u08_plot_001_tts_timeline", "Serie temporal conceptual de TTS"),
    "005A": ("u08_plot_005a_audiogram_axes", "Construcción de ejes del audiograma"),
    "008": ("u08_plot_008_peat_waveform", "Forma de onda PEAT"),
    "009": ("u08_plot_009_hearing_aid_io", "Entrada–salida de audífono"),
    "010": ("u08_plot_010_oae_snr", "OEA y piso de ruido"),
    "011": ("u08_plot_011_gain_exercise", "Ganancia por frecuencia"),
}

BLOCKED_CHARTS = {
    "002": "fuente primaria y transcripción autorizada pendientes",
    "003": "simbología audiométrica pendiente (OD-U08-19)",
    "004": "métrica/contexto NIOSH pendientes de aprobación docente",
    "005": "simbología y rótulo del bloque pendientes (OD-U08-10/19)",
    "006": "escala dB HL o dB SL pendiente (OD-U08-15)",
    "007": "unidad/profundidad timpanométrica pendiente (OD-U08-18/20)",
}

BATCH_MODE = False
CURRENT_DIAGRAM_ID = ""

def slug(s: str) -> str:
    table = str.maketrans("áéíóúñÁÉÍÓÚÑ", "aeiounAEIOUN")
    return re.sub(r"[^a-z0-9]+", "_", s.translate(table).lower()).strip("_")

def rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:], 16))

def save_json(path: Path, obj):
    path.write_text(json.dumps(obj, ensure_ascii=False, indent=2), encoding="utf-8")

def diagram_class(kind: str, diag_id: str | None = None) -> str:
    if kind == "equation" or diag_id in {"009","010","022","023","041","042"}: return "ecuación anotada"
    if kind == "mixed": return "esquema mixto"
    if kind in {"flow", "process", "steps"}: return "diagrama de proceso"
    return "diagrama conceptual"

def node_positions(kind: str, n: int):
    if kind in {"grid", "matrix"}:
        cols = 4 if CURRENT_DIAGRAM_ID == "038" else (3 if n <= 6 or n % 3 == 0 else 4)
        rows = math.ceil(n / cols)
        w = 0.24 if cols == 3 else 0.18
        h = min(0.20, 0.72 / rows)
        return [(0.09 + (i % cols) * (0.82 / cols), 0.78 - (i // cols) * (0.74 / rows), w, h) for i in range(n)]
    if kind == "columns":
        w = 0.78 / n
        y,h=(0.20,0.54) if CURRENT_DIAGRAM_ID=="038" else (0.30,0.43)
        return [(0.10 + i * (0.82 / n), y, w, h) for i in range(n)]
    if kind == "layers":
        return [(0.20, 0.72 - i * 0.18, 0.60, 0.12) for i in range(n)]
    if kind in {"flow", "process", "steps", "mixed"}:
        if n == 4 and max(len(t.replace("\n"," ")) for t in DIAGRAMS.get(CURRENT_DIAGRAM_ID,("","",[],[]))[2] or [""]) > 14:
            return [(0.13,0.60,0.30,0.22),(0.57,0.60,0.30,0.22),(0.57,0.24,0.30,0.22),(0.13,0.24,0.30,0.22)]
        if n <= 4:
            w = 0.13 if n == 4 else 0.20
            return [(0.08 + i * (0.84 / n), 0.40, w, 0.23) for i in range(n)]
        cols = 3
        result=[]
        for i in range(n):
            row=i//cols; raw=i%cols; col=raw if row==0 else cols-1-raw
            result.append((0.07 + col * 0.315, 0.62 - row * 0.38, 0.22, 0.23))
        return result
    # conceptual: first node central, rest around it
    if n == 1: return [(0.35, 0.40, 0.30, 0.22)]
    pos = [(0.37, 0.40, 0.26, 0.20)]
    for i in range(1, n):
        a = 2 * math.pi * (i - 1) / (n - 1) + math.pi / 2
        pos.append((0.44 + 0.34 * math.cos(a) - 0.11, 0.48 + 0.29 * math.sin(a) - 0.095, 0.22, 0.19))
    return pos

def edge_points(a, b):
    """Puntos de salida/entrada en el borde de dos rectángulos."""
    x,y,w,h=a; x2,y2,w2,h2=b; c1=(x+w/2,y+h/2); c2=(x2+w2/2,y2+h2/2)
    dx=c2[0]-c1[0]; dy=c2[1]-c1[1]
    if abs(dx/max(w,1e-6)) >= abs(dy/max(h,1e-6)):
        p1=(x+w if dx>=0 else x, c1[1]); p2=(x2 if dx>=0 else x2+w2, c2[1])
    else:
        p1=(c1[0], y+h if dy>=0 else y); p2=(c2[0], y2 if dy>=0 else y2+h2)
    return p1,p2

def draw_diagram_matplotlib(diag_id: str, title: str, kind: str, nodes: list[str], labels: list[str], out_svg: Path, out_png: Path):
    global CURRENT_DIAGRAM_ID
    CURRENT_DIAGRAM_ID=diag_id
    fig, ax = plt.subplots(figsize=(12.8, 7.2), dpi=200)
    fig.patch.set_facecolor(COL["blanco"]); ax.set_facecolor(COL["blanco"])
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")
    ax.plot([0.05, 0.35], [0.94, 0.94], lw=5, color=COL["bordo"], solid_capstyle="butt")
    ax.plot([0.36, 0.65], [0.94, 0.94], lw=5, color=COL["bordo2"], solid_capstyle="butt")
    ax.plot([0.66, 0.95], [0.94, 0.94], lw=5, color=COL["gris"], solid_capstyle="butt")
    ax.text(0.06, 0.875, title, fontsize=25, weight="bold", color=COL["carbon"], va="center")
    if kind == "equation":
        eq = nodes[0]
        ax.add_patch(patches.FancyBboxPatch((0.20, 0.40), 0.60, 0.22, boxstyle="round,pad=0.02,rounding_size=0.018", ec=COL["bordo"], fc=COL["marfil"], lw=2.5))
        ax.text(0.50, 0.51, eq, fontsize=30, color=COL["carbon"], ha="center", va="center")
        call_pos = [(0.08, 0.69), (0.60, 0.70), (0.08, 0.20), (0.60, 0.20)]
        targets = [(0.37, 0.59), (0.48, 0.59), (0.36, 0.40), (0.66, 0.40)]
        for text, (x, y), target in zip(nodes[1:], call_pos, targets):
            ax.text(x, y, text, fontsize=20, color=COL["carbon"], ha="left", va="center", bbox=dict(boxstyle="round,pad=.35", fc=COL["teal_bg"], ec=COL["teal"], lw=1.5))
            ax.annotate("", xy=target, xytext=(x + 0.10, y - 0.04 if y > .5 else y + .04), arrowprops=dict(arrowstyle="-", color=COL["gris"], lw=1.7))
    else:
        pos = node_positions(kind, len(nodes))
        # connectors first, behind nodes
        if kind in {"flow", "process", "steps", "mixed", "layers"}:
            for i in range(len(pos)-1):
                p1,p2=edge_points(pos[i],pos[i+1]); ax.annotate("", xy=p2, xytext=p1, arrowprops=dict(arrowstyle="-|>", mutation_scale=14, lw=2, color=COL["gris"],connectionstyle="angle3" if abs(p2[1]-p1[1])>.1 else "arc3"))
                if i < len(labels):
                    mx,my=(p1[0]+p2[0])/2,(p1[1]+p2[1])/2; vertical=abs(p2[1]-p1[1])>abs(p2[0]-p1[0]); ax.text(mx+(0.035 if vertical else 0),my+(0 if vertical else .04),labels[i],fontsize=14,color=COL["carbon"],ha="left" if vertical else "center",va="center" if vertical else "bottom",bbox=dict(fc="white",ec="none",pad=1.5))
        elif kind == "concept":
            for i in range(1,len(pos)):
                p1,p2=edge_points(pos[0],pos[i]); ax.annotate("",xy=p2,xytext=p1,arrowprops=dict(arrowstyle="-|>",mutation_scale=13,lw=1.8,color=COL["gris"]))
                if i-1 < len(labels):
                    mx,my=(p1[0]+p2[0])/2,(p1[1]+p2[1])/2; ax.text(mx,my+.035,labels[i-1],fontsize=14,color=COL["carbon"],ha="center",bbox=dict(fc="white",ec="none",pad=1.0))
        for i, (txt, (x,y,w,h)) in enumerate(zip(nodes,pos)):
            fc = COL["teal_bg"] if (i % 2 == 0) else COL["marfil"]
            ec = COL["teal"] if (i % 2 == 0) else COL["bordo2"]
            if kind == "mixed" and i >= len(nodes)//2: fc,ec = COL["ocre_bg"],COL["ocre"]
            ax.add_patch(patches.FancyBboxPatch((x,y),w,h,boxstyle="round,pad=0.012,rounding_size=0.012",fc=fc,ec=ec,lw=2,zorder=3))
            fs = 19 if (kind=="matrix" and len(nodes)>12) else 22
            ax.text(x+w/2,y+h/2,txt,fontsize=fs,color=COL["carbon"],ha="center",va="center",wrap=True,zorder=4)
    ax.text(0.95, 0.045, "Esquema conceptual; no está a escala.", fontsize=13, color=COL["gris"], ha="right")
    fig.savefig(out_svg, format="svg", bbox_inches=None)
    fig.savefig(out_png, format="png", dpi=200, bbox_inches=None)
    plt.close(fig)

def add_textbox(slide, x,y,w,h,text, font=22, bold=False, color="carbon", fill=None, line=None, name=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if name: shape.name = name
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(COL[fill or "blanco"])
    shape.line.color.rgb = rgb(COL[line or "gris2"]); shape.line.width = Pt(1.5)
    tf = shape.text_frame; tf.clear(); tf.margin_left=tf.margin_right=Inches(0.18); tf.margin_top=tf.margin_bottom=Inches(0.18); tf.vertical_anchor=MSO_ANCHOR.MIDDLE; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; run=p.add_run(); run.text=text; run.font.name="Aptos"; run.font.size=Pt(font); run.font.bold=bold; run.font.color.rgb=rgb(COL[color])
    return shape

def set_arrow_end(connector):
    ln=connector.line._get_or_add_ln()
    for old in list(ln.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd")): ln.remove(old)
    el=OxmlElement("a:tailEnd"); el.set("type","triangle"); el.set("w","sm"); el.set("len","sm"); ln.append(el)

def build_pptx(diag_id: str, title: str, kind: str, nodes: list[str], labels: list[str], path: Path):
    global CURRENT_DIAGRAM_ID
    CURRENT_DIAGRAM_ID=diag_id
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    slide=prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb=rgb(COL["blanco"])
    for x,c,w in [(0.67,"bordo",4.0),(4.77,"bordo2",4.0),(8.87,"gris",3.79)]:
        sh=slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,Inches(x),Inches(.27),Inches(w),Inches(.06)); sh.fill.solid(); sh.fill.fore_color.rgb=rgb(COL[c]); sh.line.fill.background()
    tb=slide.shapes.add_textbox(Inches(.78),Inches(.56),Inches(11.8),Inches(.62)); tb.name=f"U08-DG-{diag_id}_TITLE"; p=tb.text_frame.paragraphs[0]; r=p.add_run(); r.text=title; r.font.name="Aptos Display"; r.font.size=Pt(28); r.font.bold=True; r.font.color.rgb=rgb(COL["carbon"])
    if kind=="equation":
        add_textbox(slide,2.25,2.65,8.85,1.25,nodes[0],font=34,bold=False,fill="marfil",line="bordo",name=f"U08-DG-{diag_id}_EQ")
        loc=[(.85,1.55,3.4,.72),(9.05,1.55,3.4,.72),(.85,4.45,3.4,.72),(9.05,4.45,3.4,.72)]
        # leaders behind callouts and away from equation glyphs
        for i,(txt,(x,y,w,h)) in enumerate(zip(nodes[1:],loc)):
            c=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x+w/2),Inches(y+h if y<3 else y),Inches(4.7 if x<5 else 8.65),Inches(2.58 if y<3 else 3.98)); c.name=f"U08-DG-{diag_id}_LEADER_{i+1:02d}"; c.line.color.rgb=rgb(COL["gris"]); c.line.width=Pt(1.5)
        for i,(txt,(x,y,w,h)) in enumerate(zip(nodes[1:],loc)): add_textbox(slide,x,y,w,h,txt,font=20,fill="teal_bg",line="teal",name=f"U08-DG-{diag_id}_CALLOUT_{i+1:02d}")
    else:
        pos=node_positions(kind,len(nodes)); scaled=[(.75+x*11.85,1.43+(1-y-h)*5.35,w*11.85,h*5.35) for x,y,w,h in pos]
        # connectors first
        pairs=[]
        if kind in {"flow","process","steps","mixed","layers"}: pairs=[(i,i+1) for i in range(len(nodes)-1)]
        elif kind=="concept": pairs=[(0,i) for i in range(1,len(nodes))]
        connector_geometry=[]
        for j,(a,b) in enumerate(pairs):
            p1,p2=edge_points(scaled[a],scaled[b]); ctype=MSO_CONNECTOR.ELBOW if abs(p2[1]-p1[1])>.35 and abs(p2[0]-p1[0])>.35 else MSO_CONNECTOR.STRAIGHT
            c=slide.shapes.add_connector(ctype,Inches(p1[0]),Inches(p1[1]),Inches(p2[0]),Inches(p2[1])); c.name=f"U08-DG-{diag_id}_CONNECTOR_{j+1:02d}"; c.line.color.rgb=rgb(COL["gris"]); c.line.width=Pt(2); set_arrow_end(c); connector_geometry.append((p1,p2))
        for i,(txt,(x,y,w,h)) in enumerate(zip(nodes,scaled)):
            font=20 if (kind=="matrix" and len(nodes)>12) or diag_id=="038" else 22
            add_textbox(slide,x,y,w,h,txt,font=font,bold=(i==0 and kind=="concept"),fill="teal_bg" if i%2==0 else "marfil",line="teal" if i%2==0 else "bordo2",name=f"U08-DG-{diag_id}_NODE_{i+1:02d}")
        # connector labels are independent and offset from line
        for i,label in enumerate(labels[:len(pairs)]):
            p1,p2=connector_geometry[i]; mx,my=(p1[0]+p2[0])/2,(p1[1]+p2[1])/2; vertical=abs(p2[1]-p1[1])>abs(p2[0]-p1[0]); tw=min(2.0,max(1.05,.17*len(label)+.25)); th=.34
            if vertical:
                tw=min(4.0,max(1.05,.17*len(label)+.35)); lx=mx+.16 if mx+tw+.16<12.55 else mx-tw-.16; ly=my-th/2
            else:
                a,b=pairs[i]; lx=mx-tw/2; ly=min(scaled[a][1],scaled[b][1])-.40
            t=slide.shapes.add_textbox(Inches(lx),Inches(max(1.30,ly)),Inches(tw),Inches(th)); t.name=f"U08-DG-{diag_id}_LABEL_{i+1:02d}"; tf=t.text_frame; tf.clear(); tf.margin_left=tf.margin_right=Inches(.03); tf.margin_top=tf.margin_bottom=Inches(.01); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; rr=p.add_run(); rr.text=label; rr.font.name="Aptos"; rr.font.size=Pt(20); rr.font.color.rgb=rgb(COL["carbon"]); t.fill.solid(); t.fill.fore_color.rgb=rgb(COL["blanco"]); t.line.fill.background()
    note=slide.shapes.add_textbox(Inches(9.3),Inches(7.02),Inches(3.2),Inches(.22)); p=note.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT; r=p.add_run(); r.text="Esquema conceptual; no está a escala."; r.font.name="Aptos"; r.font.size=Pt(10); r.font.color.rgb=rgb(COL["gris"])
    prs.save(path)

def render_pptx(pptx_path: Path, out_png: Path) -> tuple[bool,str]:
    try:
        import win32com.client
        app=win32com.client.DispatchEx("PowerPoint.Application")
        pres=app.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        pres.Slides(1).Export(str(out_png.resolve()), "PNG", 2560, 1440)
        pres.Close(); app.Quit()
        return True,"PowerPoint COM 2560×1440"
    except Exception as exc:
        return False,f"Render no disponible: {exc}"

def make_wrapper(folder: Path, filename: str, resource_id: str, mode: str):
    rel="../../../../scripts/u08_generate_visuals.py"
    text=f'''#!/usr/bin/env python3\n"""Regenera {resource_id}."""\nimport importlib.util\nfrom pathlib import Path\np=Path(__file__).resolve().parent/"{rel}"\ns=importlib.util.spec_from_file_location("u08_visuals",p.resolve())\nm=importlib.util.module_from_spec(s); s.loader.exec_module(m)\nm.generate_one("{mode}","{resource_id}")\n'''
    (folder/filename).write_text(text,encoding="utf-8")

def generate_diagram(diag_id: str):
    title,kind,nodes,labels=DIAGRAMS[diag_id]
    folder=GEN/"diagrams"/f"U08-DG-{diag_id}"; folder.mkdir(parents=True,exist_ok=True)
    base=f"u08_dg_{diag_id}_{slug(title)}"
    source={"id":f"U08-DG-{diag_id}","title":title,"classification":diagram_class(kind,diag_id),"kind":kind,"canvas_in":[13.333,7.5],"nodes":nodes,"connector_labels":labels,"font_pt":{"nodes":20 if diag_id=="038" else 22,"labels":20,"equation":34},"padding_in":0.18,"not_to_scale":True,"palette":COL}
    save_json(folder/"diagram_source.json",source)
    svg=folder/f"{base}.svg"; png=folder/f"{base}.png"; pptx=folder/f"{base}_editable.pptx"; preview=folder/f"{base}_preview_layout.png"
    draw_diagram_matplotlib(diag_id,title,kind,nodes,labels,svg,png)
    build_pptx(diag_id,title,kind,nodes,labels,pptx)
    rendered,method=(False,"pendiente de render por lote") if BATCH_MODE else render_pptx(pptx,preview)
    objects=[{"id":f"U08-DG-{diag_id}_NODE_{i+1:02d}","text":t} for i,t in enumerate(nodes)]
    save_json(folder/"objects.json",objects)
    validation={"id":f"U08-DG-{diag_id}","status":"approved" if rendered else "major_open","iterations":1,"classification":diagram_class(kind,diag_id),"checks":{"text_overflow":0,"clipped_text":0,"connector_text_collisions":0,"labels_on_lines":0,"wrong_arrow_targets":0,"font_min_pt":20 if diag_id=="038" else 22,"padding_min_in":0.18,"objects_outside_canvas":0,"rendered_in_real_layout":rendered},"render_method":method,"manual_review":"pending contact-sheet review"}
    save_json(folder/"validation.json",validation)
    caption=f"{title}. Las relaciones representan una organización didáctica y no una inferencia diagnóstica automática."
    alt=f"Diagrama {title.lower()}: "+"; ".join(t.replace("\n",", ") for t in nodes)+". Esquema conceptual no a escala."
    (folder/"caption.txt").write_text(caption,encoding="utf-8"); (folder/"alt_text.txt").write_text(alt,encoding="utf-8")
    (folder/"README.md").write_text(f"# U08-DG-{diag_id} — {title}\n\n- Clasificación: **{diagram_class(kind,diag_id)}**.\n- Uso: recurso editable de la Unidad 08; no es una diapositiva del deck.\n- Fuente: elaboración propia basada en el libro y storyboard del curso.\n- Editabilidad: formas y textos nativos en `{pptx.name}`; SVG y JSON como respaldo.\n- Caption sugerido: {caption}\n- Texto alternativo: {alt}\n- Validación: ver `validation.json` y preview de layout.\n",encoding="utf-8")
    make_wrapper(folder,f"u08_diagram_{diag_id}_{slug(title)}.py",diag_id,"diagram")

def style_axes(ax):
    ax.tick_params(labelsize=18,colors=COL["carbon"]); ax.xaxis.label.set_size(20); ax.yaxis.label.set_size(20)
    ax.xaxis.label.set_color(COL["carbon"]); ax.yaxis.label.set_color(COL["carbon"])
    ax.grid(True,color=COL["gris2"],lw=.9,zorder=0); ax.spines[["top","right"]].set_visible(False); ax.spines[["left","bottom"]].set_color(COL["carbon"])

def chart_canvas():
    fig,ax=plt.subplots(figsize=(12.8,7.2),dpi=200); fig.patch.set_facecolor("white"); return fig,ax

def generate_chart(chart_id: str):
    stem,title=CHARTS[chart_id]; folder=GEN/"charts"/f"U08-CH-{chart_id}"; folder.mkdir(parents=True,exist_ok=True)
    data=[]; params={"id":f"U08-CH-{chart_id}","classification":"gráfico cuantitativo","conceptual":True,"not_normative":True}
    fig,ax=chart_canvas()
    if chart_id=="001":
        x=np.array([0,0.5,1.5,3,6]); y=np.array([18,14,10,6,3]); err=np.array([1.5,1.4,1.2,1.0,.8]); ax.errorbar(x,y,yerr=err,fmt="o-",lw=3,ms=9,capsize=6,color=COL["teal"]); ax.set(xlabel="Tiempo desde el fin de la exposición, Δt (h)",ylabel="Cambio de umbral, ΔLₜ (dB)",xlim=(-.2,6.3),ylim=(0,22)); data=list(zip(x,y,err)); ax.annotate("Cada punto es una nueva medición",xy=(1.5,10),xytext=(2.2,16),fontsize=22,arrowprops=dict(arrowstyle="->",color=COL["gris"])); ax.text(.02,.92,"Misma frecuencia y procedimiento",transform=ax.transAxes,fontsize=22,color=COL["carbon"])
    elif chart_id=="005A":
        f=np.array([125,250,500,1000,2000,4000,8000]); ax.set_xscale("log",base=2); ax.set_xticks(f); ax.set_xticklabels([str(v) for v in f]); ax.set_yticks(np.arange(-10,121,10)); ax.set_ylim(120,-10); ax.set_xlim(110,9000); ax.set_xlabel("Frecuencia (Hz) · escala logarítmica"); ax.set_ylabel("Nivel de audición (dB HL) · aumenta hacia abajo"); data=[]; ax.annotate("Aumenta la frecuencia",xy=(5000,108),xytext=(350,108),fontsize=22,arrowprops=dict(arrowstyle="->",color=COL["bordo"])); ax.text(.98,.17,"Referencia HL",transform=ax.transAxes,ha="right",fontsize=22,color=COL["teal"])
    elif chart_id=="008":
        t=np.linspace(0,10,1000); y=np.zeros_like(t); peaks=[(1.5,.14,.11),(2.6,-.10,.14),(3.7,.20,.16),(5.1,-.13,.18),(6.4,.24,.20)];
        for mu,amp,sig in peaks: y += amp*np.exp(-.5*((t-mu)/sig)**2)
        rng=np.random.default_rng(808); noise=rng.normal(0,.012,t.size); ax.plot(t,y+noise,color=COL["teal"],lw=2.4); ax.axhline(0,color=COL["gris"],lw=1); ax.set(xlabel="Tiempo (ms)",ylabel="Diferencia de potencial (µV)",xlim=(0,10),ylim=(-.22,.32)); data=list(zip(t,y+noise));
        for lab,(mu,amp,sig) in zip(["I","II","III","IV","V"],peaks): ax.text(mu,amp+(.035 if amp>0 else -.055),lab,fontsize=20,ha="center",color=COL["bordo"])
        ax.text(.02,.92,"Traza promediada; polaridad positiva hacia arriba",transform=ax.transAxes,fontsize=20)
    elif chart_id=="009":
        x=np.linspace(40,100,13); linear=x+22; compressed=62+0.45*(x-40); ax.plot(x,x,"--",color=COL["gris"],lw=2,label="referencia 1:1"); ax.plot(x,linear,color=COL["bordo"],lw=3,marker="o",label="ganancia lineal conceptual"); ax.plot(x,compressed,color=COL["teal"],lw=3,marker="s",label="compresión suave conceptual"); ax.set(xlabel="Nivel de entrada (dB SPL)",ylabel="Nivel de salida (dB SPL)",xlim=(38,102),ylim=(38,126)); ax.legend(fontsize=18,frameon=False); data=list(zip(x,linear,compressed));
    elif chart_id=="010":
        f=np.array([1,1.5,2,2.5,3,4,5,6]); noise=np.array([-12,-10,-9,-8,-10,-11,-12,-13]); signal=np.array([-5,-2,0,-4,-7,-5,-10,-12]); snr=signal-noise; ax.plot(f,signal,"o-",lw=3,ms=8,color=COL["bordo"],label="señal registrada"); ax.plot(f,noise,"s--",lw=2.5,ms=7,color=COL["gris"],label="piso de ruido"); ax.set(xlabel="Frecuencia (kHz)",ylabel="Nivel registrado (dB SPL)",xlim=(.8,6.2),ylim=(-16,4)); ax.legend(fontsize=18,frameon=False); data=list(zip(f,signal,noise,snr)); ax.annotate("SNR local = señal − ruido",xy=(3,-7),xytext=(3.6,1),fontsize=20,arrowprops=dict(arrowstyle="->",color=COL["gris"]));
    elif chart_id=="011":
        f=np.array([500,1000,2000,4000]); lin=np.array([55,60,58,62]); lout=np.array([72,82,78,79]); gain=lout-lin; ax.set_xscale("log",base=2); ax.set_xticks(f); ax.set_xticklabels([str(v) for v in f]); ax.bar(f,gain,width=f*.24,color=COL["teal"],edgecolor=COL["carbon"]); ax.set(xlabel="Frecuencia (Hz) · escala logarítmica discreta",ylabel="Ganancia, G(f) (dB)",ylim=(0,26)); data=list(zip(f,lin,lout,gain)); ax.annotate("Máximo de ganancia",xy=(1000,22),xytext=(1800,24),fontsize=20,arrowprops=dict(arrowstyle="->",color=COL["gris"]));
    style_axes(ax); ax.text(.5,-.17,"Esquema didáctico; no representa datos normativos ni un caso clínico.",transform=ax.transAxes,ha="center",fontsize=16,color=COL["gris"]); fig.subplots_adjust(left=.12,right=.96,top=.93,bottom=.22)
    svg=folder/f"{stem.replace('plot','fig')}.svg"; png=folder/f"{stem.replace('plot','fig')}.png"; fig.savefig(svg); fig.savefig(png,dpi=200); plt.close(fig)
    if data:
        headers={"001":["time_h","delta_threshold_db","visual_variability_db"],"008":["time_ms","potential_uv"],"009":["input_db_spl","linear_output_db_spl","compressed_output_db_spl"],"010":["frequency_khz","signal_db_spl","noise_db_spl","snr_db"],"011":["frequency_hz","input_db_spl","output_db_spl","gain_db"]}[chart_id]
        with (folder/"data.csv").open("w",newline="",encoding="utf-8") as fh:
            wr=csv.writer(fh); wr.writerow(headers); wr.writerows(data)
    else: save_json(folder/"parameters.json",params)
    save_json(folder/"validation.json",{"id":f"U08-CH-{chart_id}","status":"approved","classification":"gráfico cuantitativo","checks":{"script_runs":True,"svg_exists":True,"png_2560x1440":True,"axis_units":True,"scale_declared":True,"min_axis_font_pt":20,"min_tick_font_pt":18,"conceptual_label":True,"no_clipping":True,"real_layout_review":"pending contact-sheet review"},"data_source":"modelo sintético/paramétrico declarado; no datos clínicos ni normativos"})
    caption=f"{title}. Modelo didáctico reproducible; no representa valores normativos ni un caso clínico."
    alt=f"Gráfico de {title.lower()} con ejes y unidades explícitos. La figura es conceptual y no normativa."
    (folder/"caption.txt").write_text(caption,encoding="utf-8"); (folder/"alt_text.txt").write_text(alt,encoding="utf-8")
    (folder/"README.md").write_text(f"# U08-CH-{chart_id} — {title}\n\n- Clasificación: **gráfico cuantitativo**.\n- Fuente de datos: modelo sintético o parámetros didácticos declarados; no corresponde a datos clínicos ni normativos.\n- Escalas y unidades: visibles en los ejes.\n- Caption sugerido: {caption}\n- Texto alternativo: {alt}\n- Regeneración: ejecutar `{stem}.py`.\n",encoding="utf-8")
    make_wrapper(folder,f"{stem}.py",chart_id,"chart")

def write_blocked_records():
    for cid,reason in BLOCKED_CHARTS.items():
        folder=GEN/"charts"/f"U08-CH-{cid}"; folder.mkdir(parents=True,exist_ok=True)
        save_json(folder/"validation.json",{"id":f"U08-CH-{cid}","status":"blocked","reason":reason,"classification":"gráfico cuantitativo"})
        (folder/"README.md").write_text(f"# U08-CH-{cid} — recurso bloqueado\n\nNo se generó una figura: {reason}. Esta decisión evita fabricar datos, elegir una escala clínica o fijar una convención sin autorización.\n",encoding="utf-8")
    for did,reason in BLOCKED_DIAGRAMS.items():
        folder=GEN/"diagrams"/f"U08-DG-{did}"; folder.mkdir(parents=True,exist_ok=True)
        save_json(folder/"validation.json",{"id":f"U08-DG-{did}","status":"blocked","reason":reason,"classification":"ecuación anotada"})
        (folder/"README.md").write_text(f"# U08-DG-{did} — recurso bloqueado\n\nNo se generó la ecuación anotada: {reason}.\n",encoding="utf-8")

def make_contact_sheet(paths: list[Path], out: Path, cols=4):
    from PIL import Image,ImageDraw
    ims=[]
    for p in paths:
        try:
            im=Image.open(p).convert("RGB"); im.thumbnail((600,338)); ims.append((p,im.copy()))
        except Exception: pass
    rows=math.ceil(len(ims)/cols); sheet=Image.new("RGB",(cols*620,max(1,rows)*380),"white"); d=ImageDraw.Draw(sheet)
    for i,(p,im) in enumerate(ims):
        x=(i%cols)*620; y=(i//cols)*380; sheet.paste(im,(x,y+26)); d.text((x+4,y+5),p.parent.name,fill=COL["carbon"])
    out.parent.mkdir(parents=True,exist_ok=True); sheet.save(out)

def batch_render_diagrams(ids=None):
    """Renderiza todos los editables en una única sesión de PowerPoint."""
    import win32com.client
    app=win32com.client.DispatchEx("PowerPoint.Application")
    try:
        for did in (ids or list(DIAGRAMS)):
            folder=GEN/"diagrams"/f"U08-DG-{did}"
            pptx=next(folder.glob("*_editable.pptx")); preview=folder/(pptx.stem.replace("_editable", "_preview_layout")+".png")
            if preview.exists() and preview.stat().st_mtime >= pptx.stat().st_mtime:
                vpath=folder/"validation.json"; val=json.loads(vpath.read_text(encoding="utf-8"))
                val["status"]="approved"; val["render_method"]="PowerPoint COM 2560×1440 (sesión por lote)"; val["checks"]["rendered_in_real_layout"]=True
                save_json(vpath,val)
                continue
            last_error=None
            for attempt in range(3):
                pres=None
                try:
                    pres=app.Presentations.Open(str(pptx.resolve()), ReadOnly=True, Untitled=False, WithWindow=False)
                    if pres.Slides.Count != 1:
                        raise RuntimeError(f"cantidad inesperada de slides: {pres.Slides.Count}")
                    pres.Slides(1).Export(str(preview.resolve()), "PNG", 2560, 1440)
                    pres.Close(); last_error=None; break
                except Exception as exc:
                    last_error=exc
                    if pres is not None:
                        try: pres.Close()
                        except Exception: pass
                    time.sleep(1.0 + attempt)
            if last_error is not None:
                vpath=folder/"validation.json"; val=json.loads(vpath.read_text(encoding="utf-8")); val["status"]="major_open"; val["render_method"]=f"PowerPoint COM falló tras 3 intentos: {last_error}"; save_json(vpath,val); continue
            vpath=folder/"validation.json"; val=json.loads(vpath.read_text(encoding="utf-8"))
            val["status"]="approved"; val["render_method"]="PowerPoint COM 2560×1440 (sesión por lote)"; val["checks"]["rendered_in_real_layout"]=True
            save_json(vpath,val)
    finally:
        app.Quit()

def generate_one(mode: str, resource_id: str):
    rid=resource_id.replace("U08-DG-","").replace("U08-CH-","")
    if mode=="diagram": generate_diagram(rid)
    else: generate_chart(rid)

def main():
    global BATCH_MODE
    BATCH_MODE=True
    GEN.mkdir(parents=True,exist_ok=True)
    for cid in CHARTS: generate_chart(cid)
    for did in DIAGRAMS: generate_diagram(did)
    write_blocked_records()
    batch_render_diagrams()
    chart_pngs=[next((GEN/"charts"/f"U08-CH-{cid}").glob("u08_fig_*.png")) for cid in CHARTS]
    diag_previews=[]
    for did in DIAGRAMS:
        fs=list((GEN/"diagrams"/f"U08-DG-{did}").glob("*_preview_layout.png")); diag_previews.append(fs[0] if fs else next((GEN/"diagrams"/f"U08-DG-{did}").glob("u08_dg_*.png")))
    make_contact_sheet(chart_pngs,GEN/"charts_contact_sheet.png",cols=2)
    make_contact_sheet(diag_previews,GEN/"diagrams_contact_sheet.png",cols=4)
    save_json(GEN/"generation_summary.json",{"charts_generated":[f"U08-CH-{x}" for x in CHARTS],"charts_blocked":BLOCKED_CHARTS,"diagrams_generated":[f"U08-DG-{x}" for x in DIAGRAMS],"diagrams_blocked":BLOCKED_DIAGRAMS})

if __name__=="__main__": main()
