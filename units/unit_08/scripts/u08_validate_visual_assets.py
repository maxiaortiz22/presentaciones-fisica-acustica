#!/usr/bin/env python3
"""Valida integridad y editabilidad de los recursos visuales de U08."""

from __future__ import annotations

import json
import sys
import xml.etree.ElementTree as ET
from pathlib import Path

from PIL import Image

ROOT=Path(__file__).resolve().parents[3]
UNIT=ROOT/"units"/"unit_08"; GEN=UNIT/"assets"/"generated"
runtime=Path(r"C:\Users\maxia\.cache\codex-runtimes\codex-primary-runtime\dependencies\python\Lib\site-packages")
if runtime.exists(): sys.path.append(str(runtime))
from pptx import Presentation

source=UNIT/"scripts"/"u08_generate_visuals.py"; ns={"__file__":str(source),"__name__":"u08v"}
exec(compile(source.read_text(encoding="utf-8"),str(source),"exec"),ns)

errors=[]; checks=[]
def require(cond,msg):
    checks.append(msg)
    if not cond: errors.append(msg)

for cid in ns["CHARTS"]:
    folder=GEN/"charts"/f"U08-CH-{cid}"
    py=list(folder.glob("*.py")); svg=list(folder.glob("*.svg")); png=list(folder.glob("u08_fig_*.png"))
    require(len(py)==1,f"U08-CH-{cid}: un script")
    require(len(svg)==1,f"U08-CH-{cid}: un SVG")
    require(len(png)==1,f"U08-CH-{cid}: un PNG")
    require((folder/"README.md").exists() and (folder/"caption.txt").exists() and (folder/"alt_text.txt").exists() and (folder/"validation.json").exists(),f"U08-CH-{cid}: documentación completa")
    require((folder/"data.csv").exists() or (folder/"parameters.json").exists(),f"U08-CH-{cid}: datos o parámetros")
    if svg:
        try: ET.parse(svg[0]); ok=True
        except Exception: ok=False
        require(ok,f"U08-CH-{cid}: SVG parseable")
    if png: require(Image.open(png[0]).size==(2560,1440),f"U08-CH-{cid}: PNG 2560×1440")

for cid in ns["BLOCKED_CHARTS"]:
    folder=GEN/"charts"/f"U08-CH-{cid}"
    require((folder/"README.md").exists() and (folder/"validation.json").exists(),f"U08-CH-{cid}: bloqueo documentado")
    require(not list(folder.glob("*.svg")) and not list(folder.glob("*.png")),f"U08-CH-{cid}: sin figura fabricada")

for did in ns["DIAGRAMS"]:
    folder=GEN/"diagrams"/f"U08-DG-{did}"
    needed=["diagram_source.json","objects.json","README.md","caption.txt","alt_text.txt","validation.json"]
    require(all((folder/x).exists() for x in needed),f"U08-DG-{did}: metadatos completos")
    py=list(folder.glob("*.py")); svg=list(folder.glob("*.svg")); png=[p for p in folder.glob("u08_dg_*.png") if "preview_layout" not in p.name]; previews=list(folder.glob("*_preview_layout.png")); pptx=list(folder.glob("*_editable.pptx"))
    require(len(py)==1 and len(svg)==1 and len(png)==1 and len(previews)==1 and len(pptx)==1,f"U08-DG-{did}: paquete binario completo")
    if svg:
        try: ET.parse(svg[0]); ok=True
        except Exception: ok=False
        require(ok,f"U08-DG-{did}: SVG parseable")
    if png: require(Image.open(png[0]).size==(2560,1440),f"U08-DG-{did}: PNG 2560×1440")
    if previews: require(Image.open(previews[0]).size==(2560,1440),f"U08-DG-{did}: preview 2560×1440")
    if pptx:
        prs=Presentation(pptx[0]); require(len(prs.slides)==1,f"U08-DG-{did}: PPTX individual de una slide"); require(len(prs.slides[0].shapes)>=5,f"U08-DG-{did}: formas nativas editables")
    val=json.loads((folder/"validation.json").read_text(encoding="utf-8")); require(val.get("status")=="approved",f"U08-DG-{did}: validación aprobada")

for did in ns["BLOCKED_DIAGRAMS"]:
    folder=GEN/"diagrams"/f"U08-DG-{did}"
    require((folder/"README.md").exists() and (folder/"validation.json").exists(),f"U08-DG-{did}: bloqueo documentado")
    require(not list(folder.glob("*.svg")) and not list(folder.glob("*.png")) and not list(folder.glob("*.pptx")),f"U08-DG-{did}: sin visual fabricado")

result={"status":"passed" if not errors else "failed","date":"2026-08-12","checks_run":len(checks),"errors":errors,"counts":{"charts_approved":len(ns["CHARTS"]),"charts_blocked":len(ns["BLOCKED_CHARTS"]),"diagrams_approved":len(ns["DIAGRAMS"]),"diagrams_blocked":len(ns["BLOCKED_DIAGRAMS"])}}
(GEN/"integrity_validation.json").write_text(json.dumps(result,ensure_ascii=False,indent=2),encoding="utf-8")
print(json.dumps(result,ensure_ascii=False,indent=2))
raise SystemExit(1 if errors else 0)
